import os
import json
import re
from datetime import datetime
from typing import Dict, Any, List

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
except ImportError:
    Workbook = None

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
except ImportError:
    SimpleDocTemplate = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptxRGBColor
except ImportError:
    Presentation = None

try:
    from PIL import Image as PILImage, ImageDraw, ImageFont
except ImportError:
    PILImage = None

try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')
except ImportError:
    plt = None


class ExportGenerator:
    """Generate output files in various formats based on analysis results"""

    def _require_openpyxl(self):
        if Workbook is None:
            raise ImportError("openpyxl is required. Install with: pip install openpyxl")

    def _require_docx(self):
        if Document is None:
            raise ImportError("python-docx is required. Install with: pip install python-docx")

    def _require_reportlab(self):
        if SimpleDocTemplate is None:
            raise ImportError("reportlab is required. Install with: pip install reportlab")

    def _require_pptx(self):
        if Presentation is None:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    def _require_pillow(self):
        if PILImage is None:
            raise ImportError("Pillow is required. Install with: pip install Pillow")

    def generate(self, filepath: str, output_format: str, analysis: Dict, extracted_data: Dict, output_dir: str) -> str:
        """Generate output file in the specified format"""
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        base_filename = os.path.splitext(os.path.basename(filepath))[0]

        if output_format == 'pdf':
            return self._generate_pdf(filepath, analysis, extracted_data, output_dir, f"{base_filename}_{timestamp}")
        elif output_format == 'docx':
            return self._generate_docx(filepath, analysis, extracted_data, output_dir, f"{base_filename}_{timestamp}")
        elif output_format in ('xlsx', 'excel'):
            return self._generate_excel(filepath, analysis, extracted_data, output_dir, f"{base_filename}_{timestamp}")
        elif output_format in ('pptx', 'ppt'):
            return self._generate_pptx(filepath, analysis, extracted_data, output_dir, f"{base_filename}_{timestamp}")
        elif output_format == 'png':
            return self._generate_png(filepath, analysis, extracted_data, output_dir, f"{base_filename}_{timestamp}")
        else:
            raise ValueError(f"Unsupported format: {output_format}")
    
    def _prepare_content(self, analysis: Dict, extracted_data: Dict) -> Dict:
        """Prepare unified content for all export formats"""
        content = {
            'title': 'PDF Analysis Report',
            'date': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'summary': '',
            'sections': [],
            'tables': [],
            'stats': {}
        }
        
        # Get analysis result
        result = analysis.get('result', {})
        analysis_type = analysis.get('analysis_type', 'default')
        
        # Build sections based on analysis type
        if analysis_type == 'summary':
            content['title'] = 'Document Summary Report'
            content['summary'] = result.get('summary_text', '')
            content['sections'].append({
                'heading': 'Key Points',
                'items': result.get('key_points', [])
            })
            content['sections'].append({
                'heading': 'Document Structure',
                'items': result.get('document_structure', [])
            })
            content['stats'] = result.get('statistics', {})
        
        elif analysis_type == 'extract':
            content['title'] = 'Extracted Data Report'
            extracted = result.get('extracted_data', {})
            for key, values in extracted.items():
                if isinstance(values, list) and values:
                    content['sections'].append({
                        'heading': key.replace('_', ' ').title(),
                        'items': [str(v) for v in values[:50]]
                    })
        
        elif analysis_type == 'find':
            content['title'] = 'Search Results Report'
            results = result.get('results', {})
            for keyword, data in results.items():
                sentences = data.get('sentences', [])
                if sentences:
                    content['sections'].append({
                        'heading': f'Results for "{keyword}" ({data.get("total_occurrences", 0)} occurrences)',
                        'items': sentences[:10]
                    })
        
        elif analysis_type == 'calculate':
            content['title'] = 'Statistical Analysis Report'
            stats = {k: v for k, v in result.items() if k != 'operation' and k != 'table_data'}
            content['sections'].append({
                'heading': 'Statistics',
                'items': [f"{k.replace('_', ' ').title()}: {v}" for k, v in stats.items() if isinstance(v, (int, float, str))]
            })
            content['stats'] = stats
        
        elif analysis_type == 'compare':
            content['title'] = 'Comparison Report'
            comparison_data = result.get('comparison_data', [])
            for item in comparison_data:
                content['sections'].append({
                    'heading': f'Comparison Item {item.get("table_index", item.get("page", ""))}',
                    'items': [f"{k}: {v}" for k, v in item.items()]
                })
        
        elif analysis_type == 'organize':
            content['title'] = 'Organized Data Report'
            sections = result.get('sections', [])
            for section in sections:
                content['sections'].append({
                    'heading': section.get('title', 'Section'),
                    'items': [f"Subsection: {s.get('title', '')} (Page {s.get('page', '')})" for s in section.get('subsections', [])]
                })
            
            tables = result.get('tables', [])
            for t in tables:
                content['tables'].append({
                    'headers': t.get('headers', []),
                    'data': t.get('data', [])
                })
        
        else:
            content['title'] = 'PDF Analysis Report'
            content['summary'] = result.get('analysis', 'Document analysis results')
            key_sentences = result.get('key_sentences', [])
            if key_sentences:
                content['sections'].append({
                    'heading': 'Key Content',
                    'items': key_sentences[:20]
                })
            
            # Add extracted tables
            tables = extracted_data.get('tables', [])
            for t in tables[:5]:
                content['tables'].append({
                    'headers': t.get('headers', []),
                    'data': t.get('rows', [])[:10]
                })
            
            # Add text stats
            text_data = extracted_data.get('text', {})
            content['stats'] = {
                'Pages': text_data.get('pages', 'N/A'),
                'Paragraphs': text_data.get('paragraph_count', 'N/A'),
                'Sentences': text_data.get('sentence_count', 'N/A'),
                'Words': text_data.get('word_count', 'N/A'),
                'Tables Found': len(extracted_data.get('tables', [])),
                'Images Found': extracted_data.get('images', {}).get('total_images', 'N/A')
            }
        
        return content
    
    def _generate_pdf(self, filepath: str, analysis: Dict, extracted_data: Dict, output_dir: str, filename: str) -> str:
        """Generate PDF report"""
        self._require_reportlab()
        output_path = os.path.join(output_dir, f"{filename}.pdf")
        content = self._prepare_content(analysis, extracted_data)
        
        doc = SimpleDocTemplate(
            output_path,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c5282'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['BodyText'],
            fontSize=11,
            leading=16,
            alignment=TA_JUSTIFY
        )
        
        meta_style = ParagraphStyle(
            'MetaStyle',
            parent=styles['Normal'],
            fontSize=10,
            textColor=colors.gray,
            alignment=TA_CENTER
        )
        
        story = []
        
        # Title
        story.append(Paragraph(content['title'], title_style))
        story.append(Paragraph(f"Generated on: {content['date']}", meta_style))
        story.append(Spacer(1, 0.3 * inch))
        
        # Summary
        if content['summary']:
            story.append(Paragraph("Summary", heading_style))
            story.append(Paragraph(content['summary'], body_style))
            story.append(Spacer(1, 0.2 * inch))
        
        # Statistics
        if content['stats']:
            story.append(Paragraph("Document Statistics", heading_style))
            stats_data = [["Metric", "Value"]]
            for key, value in content['stats'].items():
                if isinstance(value, (int, float, str)):
                    stats_data.append([str(key), str(value)])
            
            if len(stats_data) > 1:
                stats_table = Table(stats_data, colWidths=[3 * inch, 3 * inch])
                stats_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5282')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f7fafc')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e2e8f0')),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 10),
                    ('TOPPADDING', (0, 1), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
                ]))
                story.append(stats_table)
                story.append(Spacer(1, 0.2 * inch))
        
        # Sections
        for section in content['sections']:
            story.append(Paragraph(section['heading'], heading_style))
            for item in section['items']:
                # Escape XML special characters
                item = str(item).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(f"• {item}", body_style))
            story.append(Spacer(1, 0.2 * inch))
        
        # Tables
        for table in content['tables']:
            if table['headers'] and table['data']:
                story.append(Paragraph("Extracted Table", heading_style))
                table_data = [table['headers']] + table['data']
                
                # Convert to strings
                table_data = [[str(cell) for cell in row] for row in table_data]
                
                col_width = 6 * inch / len(table['headers'])
                pdf_table = Table(table_data, colWidths=[col_width] * len(table['headers']))
                pdf_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5282')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f7fafc')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e2e8f0')),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                    ('TOPPADDING', (0, 1), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                story.append(pdf_table)
                story.append(Spacer(1, 0.2 * inch))
        
        doc.build(story)
        return output_path
    
    def _generate_docx(self, filepath: str, analysis: Dict, extracted_data: Dict, output_dir: str, filename: str) -> str:
        """Generate DOCX report"""
        self._require_docx()
        output_path = os.path.join(output_dir, f"{filename}.docx")
        content = self._prepare_content(analysis, extracted_data)
        
        doc = Document()
        
        # Title
        title = doc.add_heading(content['title'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Date
        date_para = doc.add_paragraph(f"Generated on: {content['date']}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
        doc.add_paragraph()
        
        # Summary
        if content['summary']:
            doc.add_heading('Summary', level=1)
            doc.add_paragraph(content['summary'])
            doc.add_paragraph()
        
        # Statistics
        if content['stats']:
            doc.add_heading('Document Statistics', level=1)
            for key, value in content['stats'].items():
                if isinstance(value, (int, float, str)):
                    p = doc.add_paragraph(style='List Bullet')
                    p.add_run(f"{key}: ").bold = True
                    p.add_run(str(value))
            doc.add_paragraph()
        
        # Sections
        for section in content['sections']:
            doc.add_heading(section['heading'], level=1)
            for item in section['items']:
                doc.add_paragraph(str(item), style='List Bullet')
            doc.add_paragraph()
        
        # Tables
        for table_data in content['tables']:
            if table_data['headers'] and table_data['data']:
                doc.add_heading('Extracted Table', level=1)
                table = doc.add_table(rows=1, cols=len(table_data['headers']))
                table.style = 'Table Grid'
                
                # Header row
                hdr_cells = table.rows[0].cells
                for i, header in enumerate(table_data['headers']):
                    hdr_cells[i].text = str(header)
                    hdr_cells[i].paragraphs[0].runs[0].font.bold = True
                    hdr_cells[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                    hdr_cells[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                    # Color the cell background (blue)
                    from docx.oxml.ns import qn
                    from docx.oxml import OxmlElement
                    shading = OxmlElement('w:shd')
                    shading.set(qn('w:fill'), '2c5282')
                    hdr_cells[i]._tc.get_or_add_tcPr().append(shading)
                
                # Data rows
                for row_data in table_data['data']:
                    row_cells = table.add_row().cells
                    for i, cell in enumerate(row_data):
                        row_cells[i].text = str(cell)
                
                doc.add_paragraph()
        
        doc.save(output_path)
        return output_path
    
    def _generate_excel(self, filepath: str, analysis: Dict, extracted_data: Dict, output_dir: str, filename: str) -> str:
        """Generate Excel report"""
        self._require_openpyxl()
        output_path = os.path.join(output_dir, f"{filename}.xlsx")
        content = self._prepare_content(analysis, extracted_data)
        
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Summary sheet
        ws_summary = wb.create_sheet("Summary")
        ws_summary['A1'] = content['title']
        ws_summary['A1'].font = Font(size=18, bold=True, color='1a365d')
        ws_summary['A2'] = f"Generated: {content['date']}"
        ws_summary['A2'].font = Font(size=10, color='808080')
        
        if content['summary']:
            ws_summary['A4'] = "Summary:"
            ws_summary['A4'].font = Font(bold=True, size=12)
            ws_summary['A5'] = content['summary']
            ws_summary['A5'].alignment = Alignment(wrap_text=True)
        
        # Statistics sheet
        if content['stats']:
            ws_stats = wb.create_sheet("Statistics")
            ws_stats['A1'] = "Metric"
            ws_stats['B1'] = "Value"
            
            # Style header
            for cell in ws_stats[1]:
                cell.font = Font(bold=True, color='FFFFFF')
                cell.fill = PatternFill(start_color='2c5282', end_color='2c5282', fill_type='solid')
                cell.alignment = Alignment(horizontal='center')
            
            row = 2
            for key, value in content['stats'].items():
                if isinstance(value, (int, float, str)):
                    ws_stats[f'A{row}'] = str(key)
                    ws_stats[f'B{row}'] = str(value)
                    row += 1
            
            # Auto-adjust columns
            ws_stats.column_dimensions['A'].width = 30
            ws_stats.column_dimensions['B'].width = 30
        
        # Sections sheets
        for i, section in enumerate(content['sections']):
            sheet_name = section['heading'][:31]  # Excel sheet name max 31 chars
            ws = wb.create_sheet(sheet_name)
            
            ws['A1'] = section['heading']
            ws['A1'].font = Font(bold=True, size=14, color='2c5282')
            
            row = 3
            for item in section['items']:
                ws[f'A{row}'] = str(item)
                ws[f'A{row}'].alignment = Alignment(wrap_text=True, vertical='top')
                row += 1
            
            ws.column_dimensions['A'].width = 80
        
        # Tables sheet
        for i, table in enumerate(content['tables']):
            if table['headers'] and table['data']:
                sheet_name = f"Table_{i+1}"
                ws = wb.create_sheet(sheet_name)
                
                # Headers
                for j, header in enumerate(table['headers']):
                    cell = ws.cell(row=1, column=j+1, value=str(header))
                    cell.font = Font(bold=True, color='FFFFFF')
                    cell.fill = PatternFill(start_color='2c5282', end_color='2c5282', fill_type='solid')
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                
                # Data
                for r, row_data in enumerate(table['data']):
                    for c, cell_value in enumerate(row_data):
                        cell = ws.cell(row=r+2, column=c+1, value=str(cell_value))
                        cell.alignment = Alignment(vertical='top', wrap_text=True)
                
                # Auto-adjust columns
                for col in ws.columns:
                    max_length = 0
                    column = col[0].column_letter
                    for cell in col:
                        try:
                            if cell.value:
                                max_length = max(max_length, len(str(cell.value)))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    ws.column_dimensions[column].width = adjusted_width
        
        wb.save(output_path)
        return output_path
    
    def _generate_pptx(self, filepath: str, analysis: Dict, extracted_data: Dict, output_dir: str, filename: str) -> str:
        """Generate PowerPoint presentation"""
        self._require_pptx()
        output_path = os.path.join(output_dir, f"{filename}.pptx")
        content = self._prepare_content(analysis, extracted_data)
        
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content['title']
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(26, 54, 93)
        
        subtitle.text = f"Generated: {content['date']}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(128, 128, 128)
        
        # Summary slide
        if content['summary']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Summary"
            title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(44, 82, 130)
            body.text = content['summary']
        
        # Statistics slide
        if content['stats']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Document Statistics"
            title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(44, 82, 130)
            
            tf = body.text_frame
            tf.clear()
            for key, value in content['stats'].items():
                if isinstance(value, (int, float, str)):
                    p = tf.add_paragraph()
                    p.text = f"{key}: {value}"
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)
        
        # Section slides
        for section in content['sections']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = section['heading']
            title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(44, 82, 130)
            
            tf = body.text_frame
            tf.clear()
            for item in section['items'][:10]:  # Limit items per slide
                p = tf.add_paragraph()
                p.text = f"• {str(item)[:100]}"
                p.font.size = Pt(16)
                p.space_after = Pt(8)
                p.level = 0
        
        # Table slides
        for table in content['tables']:
            if table['headers'] and table['data']:
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Extracted Table"
                title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(44, 82, 130)
                
                rows = len(table['data']) + 1
                cols = len(table['headers'])
                
                left = PptxInches(0.5)
                top = PptxInches(1.5)
                width = PptxInches(12)
                height = PptxInches(5)
                
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                
                # Fill headers
                for i, header in enumerate(table['headers']):
                    cell = table_shape.table.cell(0, i)
                    cell.text = str(header)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PptxRGBColor(44, 82, 130)
                
                # Fill data
                for r, row_data in enumerate(table['data']):
                    for c, cell_value in enumerate(row_data):
                        cell = table_shape.table.cell(r + 1, c)
                        cell.text = str(cell_value)
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        prs.save(output_path)
        return output_path
    
    def _generate_png(self, filepath: str, analysis: Dict, extracted_data: Dict, output_dir: str, filename: str) -> str:
        """Generate PNG infographic-style report"""
        self._require_pillow()
        output_path = os.path.join(output_dir, f"{filename}.png")
        content = self._prepare_content(analysis, extracted_data)
        
        # Create a large image
        img_width = 1200
        img_height = 1600
        img = PILImage.new('RGB', (img_width, img_height), color='#ffffff')
        draw = ImageDraw.Draw(img)
        
        try:
            # Try to use a nice font
            font_title = ImageFont.truetype("arial.ttf", 40)
            font_heading = ImageFont.truetype("arial.ttf", 28)
            font_body = ImageFont.truetype("arial.ttf", 18)
            font_meta = ImageFont.truetype("arial.ttf", 16)
        except:
            # Fallback to default font
            font_title = ImageFont.load_default()
            font_heading = ImageFont.load_default()
            font_body = ImageFont.load_default()
            font_meta = ImageFont.load_default()
        
        y_pos = 30
        
        # Title with background
        draw.rectangle([0, 0, img_width, 100], fill='#1a365d')
        draw.text((img_width//2, 55), content['title'], fill='#ffffff', font=font_title, anchor='mm')
        
        y_pos = 120
        
        # Date
        draw.text((img_width//2, y_pos), f"Generated: {content['date']}", fill='#808080', font=font_meta, anchor='mm')
        y_pos += 40
        
        # Summary
        if content['summary']:
            draw.rectangle([30, y_pos, img_width-30, y_pos+50], fill='#e2e8f0')
            draw.text((50, y_pos+15), "Summary", fill='#2c5282', font=font_heading)
            y_pos += 60
            
            # Wrap text
            summary = content['summary'][:500]
            words = summary.split()
            lines = []
            current_line = []
            for word in words:
                test_line = ' '.join(current_line + [word])
                bbox = draw.textbbox((0, 0), test_line, font=font_body)
                if bbox[2] - bbox[0] < img_width - 100:
                    current_line.append(word)
                else:
                    lines.append(' '.join(current_line))
                    current_line = [word]
            if current_line:
                lines.append(' '.join(current_line))
            
            for line in lines:
                draw.text((50, y_pos), line, fill='#333333', font=font_body)
                y_pos += 25
            
            y_pos += 20
        
        # Statistics
        if content['stats']:
            draw.rectangle([30, y_pos, img_width-30, y_pos+50], fill='#e2e8f0')
            draw.text((50, y_pos+15), "Statistics", fill='#2c5282', font=font_heading)
            y_pos += 60
            
            for key, value in content['stats'].items():
                if isinstance(value, (int, float, str)):
                    draw.text((50, y_pos), f"• {key}: {value}", fill='#333333', font=font_body)
                    y_pos += 25
            
            y_pos += 20
        
        # Sections
        for section in content['sections']:
            if y_pos > img_height - 200:
                break
            
            draw.rectangle([30, y_pos, img_width-30, y_pos+50], fill='#e2e8f0')
            draw.text((50, y_pos+15), section['heading'][:50], fill='#2c5282', font=font_heading)
            y_pos += 60
            
            for item in section['items'][:10]:
                if y_pos > img_height - 50:
                    break
                item_text = f"• {str(item)[:100]}"
                draw.text((50, y_pos), item_text, fill='#333333', font=font_body)
                y_pos += 25
            
            y_pos += 20
        
        # Crop if needed
        if y_pos < img_height:
            img = img.crop((0, 0, img_width, y_pos + 50))
        
        img.save(output_path, 'PNG')
        return output_path
